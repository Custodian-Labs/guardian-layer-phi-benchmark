"""Guardian Layer deck — plain-language version for a non-technical audience.
Same clinical-paper visual identity; everyday wording, no ML jargon.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GROUND=RGBColor(0xE9,0xEC,0xEA); PAPER=RGBColor(0xF3,0xF5,0xF3); INK=RGBColor(0x18,0x21,0x1E)
MUTED=RGBColor(0x5C,0x6B,0x64); GREEN=RGBColor(0x0F,0x6E,0x5C); CLAY=RGBColor(0xBD,0x5B,0x36)
LINE=RGBColor(0xC5,0xCD,0xC8); WHITE=RGBColor(0xFF,0xFF,0xFF)
SERIF="Georgia"; MONO="Consolas"
W,H=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]

def slide(bg=GROUND):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=bg
    r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0; return tb,tf
def para(tf,t,sz,c=INK,fn=SERIF,b=False,it=False,al=PP_ALIGN.LEFT,sa=6,sb=0,first=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=al; p.space_after=Pt(sa); p.space_before=Pt(sb)
    r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.name=fn; f.bold=b; f.italic=it; f.color.rgb=c
    return p,r
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(1,x,y,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def eyebrow(s,t,x=Inches(0.7),y=Inches(0.55)):
    _,tf=box(s,x,y,Inches(11.9),Inches(0.4)); para(tf,t.upper(),12,GREEN,MONO,b=True,first=True)
def rule(s,x,y,w,c=INK,h=Pt(2.2)): rect(s,x,y,w,h,fill=c)

# ===== 1 · TITLE =====
s=slide()
eyebrow(s,"Custodian Labs · Guardian Layer")
rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(2.2),Inches(11.9),Inches(2.8))
para(tf,"Does our privacy swap",46,INK,SERIF,b=True,first=True,sa=2)
para(tf,"keep the data usable?",46,GREEN,SERIF,b=True,it=True,sa=18)
para(tf,"We hide personal details by swapping them for realistic stand-ins. "
        "This is the test of whether anything downstream breaks.",20,MUTED,SERIF)
rule(s,Inches(0.7),Inches(6.5),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(6.62),Inches(11.9),Inches(0.5))
para(tf,"Tested on 1,750 real documents across 7 datasets and 7 languages.",13,MUTED,MONO,first=True)

# ===== 2 · WHAT WE DO (the swap, intuitive) =====
s=slide()
eyebrow(s,"01 · What the tool does")
rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(1.4),Inches(11.9),Inches(0.9))
para(tf,"It doesn’t black out personal info — it swaps it for a realistic fake of the same kind.",
     24,INK,SERIF,b=True,first=True)
# before card
rect(s,Inches(0.7),Inches(2.7),Inches(11.9),Inches(2.7),fill=PAPER,line=LINE)
_,tf=box(s,Inches(1.0),Inches(2.95),Inches(11.3),Inches(2.4))
para(tf,"BEFORE",11,MUTED,MONO,b=True,first=True,sa=6)
para(tf,"…a 34-year-old female treated by Anna S. at Methodist Hospital on April 12, 2023.",
     17,CLAY,MONO,sa=14)
para(tf,"AFTER",11,MUTED,MONO,b=True,sa=6)
para(tf,"…a 35-year-old female treated by Maria S. at Methodist Hospital on March 13, 2021.",
     17,GREEN,MONO,sa=14)
para(tf,"Real name → fake name · real date → fake date. The sentence still reads normally, "
        "so anything that processed the original still works on the swapped version.",13,MUTED,SERIF)

# ===== 3 · THE ANSWER (big, plain) =====
s=slide()
eyebrow(s,"02 · The result")
rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(2.0),Inches(5.4),Inches(3.0))
para(tf,"97–100%",80,GREEN,MONO,b=True,first=True,sa=8)
para(tf,"OF THE TIME, THE SWAPPED INFO IS STILL FOUND",13,MUTED,MONO,b=True)
_,tf=box(s,Inches(6.4),Inches(2.1),Inches(6.2),Inches(3.2))
para(tf,"The swap doesn’t hide the data from other tools.",24,INK,SERIF,b=True,first=True,sa=14)
para(tf,"After swapping, other AI tools still spot the sensitive details — in the same place, "
        "almost every time. We checked this across 11 different tools (GPT-5, Gemma, Qwen, Llama 8B/70B, DeepSeek, Presidio, OBI).",18,INK,SERIF,sa=12)
para(tf,"A statistical equivalence test confirms the change is too small to matter — within "
        "±2 percentage points of no change at all.",16,GREEN,SERIF,sa=12)
para(tf,"And the structure is untouched: medical notes, forms, even raw data files keep their "
        "exact format. Nothing downstream breaks.",15,MUTED,SERIF)

# ===== 4 · IT WORKS EVERYWHERE (examples) =====
s=slide()
eyebrow(s,"03 · It holds across formats & languages")
rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
def card(y,title,before,after,note):
    rect(s,Inches(0.7),y,Inches(11.9),Inches(1.65),fill=PAPER,line=LINE)
    _,tf=box(s,Inches(0.95),y+Inches(0.13),Inches(11.4),Inches(1.5))
    para(tf,title.upper(),11,MUTED,MONO,b=True,first=True,sa=5)
    para(tf,"before   "+before,12.5,CLAY,MONO,sa=3)
    para(tf,"after    "+after,12.5,GREEN,MONO,sa=4)
    para(tf,note,11,MUTED,SERIF)
card(Inches(1.35),"Doctor’s shorthand note",
     "70yo M, seen by Dr. John L. at Mt. Sinai on Feb 21, 2023",
     "73 yo M, seen by Dr. James L. at Mt. Egypt on Nov 19, 2021",
     "Even messy clinical abbreviations stay intact — only the personal details change.")
card(Inches(3.25),"German record",
     '…ermächtigen hiermit Monsignore … Mit Datum 23/07/2011',
     '…ermächtigen hiermit Fulgenzio … Mit Datum 24/07/2011',
     "Works the same across languages.")
card(Inches(5.15),"A raw data file (JSON)",
     '{ "Date": "20/05/2022", "City": "Saint-Priest", "User": "phprosdocimo" }',
     '{ "Date": "21/05/2023", "City": "Saint-Priest", "User": "phprosdocimo" }',
     "The file format is preserved exactly — only the private value is swapped, so systems that read it never break.")

# ===== 5 · ONE HONEST NOTE =====
s=slide()
eyebrow(s,"04 · One honest note")
rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(1.6),Inches(11.9),Inches(1.4))
para(tf,"How much it swaps depends on the data.",26,INK,SERIF,b=True,first=True,sa=12)
para(tf,"The tool is tuned for genuinely private information. On real medical records it swaps "
        "most of it; on general text it swaps less, because many names there aren’t actually private.",
     18,MUTED,SERIF)
rows=[("Real medical records","~80% of personal details swapped",GREEN),
      ("General / encyclopedic text","less — those names aren’t private",CLAY)]
top=Inches(3.7)
for i,(a,b,c) in enumerate(rows):
    y=top+Inches(0.75)*i
    _,t1=box(s,Inches(0.9),y,Inches(6.0),Inches(0.6)); para(t1,a,18,INK,SERIF,b=True,first=True)
    _,t2=box(s,Inches(6.7),y,Inches(6.0),Inches(0.6)); para(t2,b,16,c,SERIF,first=True)
_,tf=box(s,Inches(0.7),Inches(5.7),Inches(11.9),Inches(1.0))
para(tf,"Important: the “still found 97–100%” result is measured only on the details it does swap — "
        "so the two points are kept separate and don’t prop each other up.",14,MUTED,SERIF,first=True)

# ===== 6 · BOTTOM LINE =====
s=slide(bg=INK)
_,tf=box(s,Inches(0.9),Inches(2.4),Inches(11.5),Inches(2.6))
para(tf,"Bottom line",13,GREEN,MONO,b=True,first=True,sa=14)
para(tf,"The swap protects privacy without breaking the data.",32,WHITE,SERIF,b=True,sa=16)
para(tf,"Sensitive details are replaced by realistic stand-ins, the format is preserved exactly, "
        "and downstream tools keep working as before.",19,RGBColor(0xC5,0xCD,0xC8),SERIF)
_,tf=box(s,Inches(0.9),Inches(6.0),Inches(11.5),Inches(0.6))
para(tf,"See it live → custodianai.pages.dev",18,GREEN,MONO,b=True,first=True)

import os; os.makedirs("docs",exist_ok=True)
prs.save("docs/guardian_layer_deck.pptx")
print("saved docs/guardian_layer_deck.pptx ·",len(prs.slides._sldIdLst),"slides")
