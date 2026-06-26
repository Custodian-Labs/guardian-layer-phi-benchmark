"""Internal complete deck — plain language, includes the coverage gap.
Same clinical-paper look as the external deck. Everyday words throughout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

GROUND=RGBColor(0xE9,0xEC,0xEA); PAPER=RGBColor(0xF3,0xF5,0xF3); INK=RGBColor(0x18,0x21,0x1E)
MUTED=RGBColor(0x5C,0x6B,0x64); GREEN=RGBColor(0x0F,0x6E,0x5C); CLAY=RGBColor(0xBD,0x5B,0x36)
LINE=RGBColor(0xC5,0xCD,0xC8); WHITE=RGBColor(0xFF,0xFF,0xFF)
SERIF="Georgia"; MONO="Consolas"
W,H=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BLANK=prs.slide_layouts[6]
def slide(bg=GROUND):
    s=prs.slides.add_slide(BLANK); r=s.shapes.add_shape(1,0,0,W,H)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0; return tb,tf
def para(tf,t,sz,c=INK,fn=SERIF,b=False,it=False,al=PP_ALIGN.LEFT,sa=6,sb=0,first=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=al; p.space_after=Pt(sa); p.space_before=Pt(sb)
    r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.name=fn; f.bold=b; f.italic=it; f.color.rgb=c
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0):
    sp=s.shapes.add_shape(1,x,y,w,h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def eyebrow(s,t):
    _,tf=box(s,Inches(0.7),Inches(0.55),Inches(11.9),Inches(0.4)); para(tf,t.upper(),12,GREEN,MONO,b=True,first=True)
def rule(s,x,y,w,c=INK,h=Pt(2.2)): rect(s,x,y,w,h,fill=c)
def bar(s,x,y,wfull,frac,col):  # simple horizontal bar
    rect(s,x,y,wfull,Inches(0.32),fill=RGBColor(0xD8,0xDE,0xDA))
    rect(s,x,y,Inches(wfull.inches*frac),Inches(0.32),fill=col)

# 1 · TITLE
s=slide(); eyebrow(s,"Custodian Labs · Guardian Layer · Internal"); rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(2.3),Inches(11.9),Inches(2.6))
para(tf,"What we learned about",44,INK,SERIF,b=True,first=True,sa=2)
para(tf,"the privacy swap.",44,GREEN,SERIF,b=True,it=True,sa=18)
para(tf,"Two things matter: does it keep data usable, and does it hide enough? "
        "Here's the honest read on both.",20,MUTED,SERIF)
rule(s,Inches(0.7),Inches(6.5),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(6.62),Inches(11.9),Inches(0.5))
para(tf,"Internal — for the team. The client-facing deck covers the usability win only.",12,MUTED,MONO,first=True)

# 2 · THE GOOD NEWS (usability)
s=slide(); eyebrow(s,"01 · The good news"); rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(2.0),Inches(5.4),Inches(3.0))
para(tf,"97–100%",78,GREEN,MONO,b=True,first=True,sa=8)
para(tf,"WHEN IT SWAPS SOMETHING, OTHER TOOLS STILL FIND IT",12,MUTED,MONO,b=True)
_,tf=box(s,Inches(6.4),Inches(2.1),Inches(6.2),Inches(3.2))
para(tf,"The swap keeps data usable.",24,INK,SERIF,b=True,first=True,sa=14)
para(tf,"After swapping, other tools still spot the sensitive details almost every time, "
        "and the format (medical notes, forms, data files) stays exactly the same. "
        "Nothing downstream breaks.",18,INK,SERIF,sa=10)
para(tf,"This is the client-facing win — solid across all 11 tools and 7 datasets.",15,MUTED,SERIF)

# 3 · THE GAP (coverage, simple)
s=slide(); eyebrow(s,"02 · The honest gap — how much does it hide?"); rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(1.35),Inches(11.9),Inches(0.7))
para(tf,"On real medical records, it hides names and dates well — but misses about half of "
        "ID numbers and addresses.",19,INK,SERIF,b=True,first=True)
rows=[("Names","hides ~3 of 4",0.75,GREEN),
      ("Dates / ages","hides ~3 of 4",0.73,GREEN),
      ("ID numbers (patient ID, phone)","hides ~1 of 2",0.51,CLAY),
      ("Addresses / places","hides ~1 of 2",0.44,CLAY)]
top=Inches(2.5)
for i,(name,lab,frac,col) in enumerate(rows):
    y=top+Inches(0.92)*i
    _,t1=box(s,Inches(0.9),y,Inches(4.6),Inches(0.4)); para(t1,name,16,INK,SERIF,b=True,first=True)
    bar(s,Inches(5.6),y+Inches(0.02),Inches(4.6),frac,col)
    _,t2=box(s,Inches(10.5),y,Inches(2.4),Inches(0.4)); para(t2,lab,14,col,MONO,b=True,first=True)

# 4 · WHY THE GAP IS FIXABLE
s=slide(); eyebrow(s,"03 · The good news about the gap"); rule(s,Inches(0.7),Inches(1.0),Inches(11.9))
_,tf=box(s,Inches(0.7),Inches(2.0),Inches(5.4),Inches(3.0))
para(tf,"~3 in 4",78,GREEN,MONO,b=True,first=True,sa=8)
para(tf,"OF THE MISSES, THE TOOL ALREADY KNEW IT WAS PRIVATE",12,MUTED,MONO,b=True)
_,tf=box(s,Inches(6.4),Inches(2.1),Inches(6.2),Inches(3.2))
para(tf,"It's a swapping gap, not a finding gap.",24,INK,SERIF,b=True,first=True,sa=14)
para(tf,"For most of what slips through, the tool had already flagged it as sensitive — "
        "it just didn't swap it. So the fix is in the swap step, not in teaching it what's "
        "private.",18,INK,SERIF,sa=10)
para(tf,"That's the easier kind of problem to fix.",15,MUTED,SERIF)

# 5 · BOTTOM LINE (balanced)
s=slide(bg=INK)
_,tf=box(s,Inches(0.9),Inches(2.1),Inches(11.5),Inches(3.4))
para(tf,"Bottom line",13,GREEN,MONO,b=True,first=True,sa=14)
para(tf,"Great at keeping data usable. Next: hide more completely.",30,WHITE,SERIF,b=True,sa=18)
para(tf,"• The swap doesn't break anything — that's ready to show clients.",18,RGBColor(0xC5,0xCD,0xC8),SERIF,sa=8)
para(tf,"• It hides names and dates well, but misses ~half of ID numbers and addresses.",18,RGBColor(0xC5,0xCD,0xC8),SERIF,sa=8)
para(tf,"• Most misses are already detected — fixing the swap step should close the gap.",18,RGBColor(0xC5,0xCD,0xC8),SERIF)
_,tf=box(s,Inches(0.9),Inches(6.2),Inches(11.5),Inches(0.5))
para(tf,"Full data → custodianai.pages.dev",16,GREEN,MONO,b=True,first=True)

import os; os.makedirs("docs",exist_ok=True)
prs.save("docs/guardian_layer_internal.pptx")
print("saved docs/guardian_layer_internal.pptx ·",len(prs.slides._sldIdLst),"slides")
